"""M365 File Reader — Dateien aus SharePoint/OneDrive ueber Graph API lesen.

Usage:
    python m365_file_reader.py search "Dateiname oder Suchbegriff"
    python m365_file_reader.py read URL_ODER_PFAD
    python m365_file_reader.py read URL_ODER_PFAD --download ZIEL_PFAD

Befehle:
    search   Datei via Graph Search API finden (liefert driveId, itemId, Name, URL)
    read     Datei-Inhalt lesen und als Text ausgeben
             --download PFAD  Optional: Datei zusaetzlich lokal speichern

Token:
    Nutzt den gleichen Token-Cache wie copilot_search.py.
    Bei abgelaufenem Token: Exit-Code 2 (Agent holt neuen via NAA).

Unterstuetzte Formate:
    .pptx  — Text pro Folie (inkl. Tabellen, Notes)
    .xlsx  — Alle Sheets als CSV (Semikolon-getrennt)
    .docx  — Volltext mit Ueberschriften
    .pdf   — Text-Extraktion (falls pdfplumber installiert)
    .png/.jpg/.jpeg/.gif/.bmp/.tiff/.svg — Bild-Download + Metadaten
    .csv   — Erste 200 Zeilen als Markdown-Tabelle
    .txt/.md/.json/.xml — Plaintext
    andere — Nur Metadaten + optionaler Download
"""

import csv
import json
import os
import sys
import time
from io import BytesIO, StringIO
from pathlib import Path
from urllib.parse import unquote, urlparse

import requests

# Windows-Konsolen-Encoding auf UTF-8 setzen
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    os.environ.setdefault("PYTHONIOENCODING", "utf-8")

CACHE_FILE = Path(__file__).resolve().parent.parent / "userdata" / "tmp" / ".graph_token_cache.json"
GRAPH_BASE = "https://graph.microsoft.com/v1.0"
MIN_TOKEN_LIFETIME = 120


# ---------------------------------------------------------------------------
# Token helpers (shared logic with copilot_search.py)
# ---------------------------------------------------------------------------

def _load_cached_token() -> str | None:
    if not CACHE_FILE.exists():
        return None
    try:
        with open(CACHE_FILE) as f:
            cache = json.load(f)
        if cache.get("exp", 0) > time.time() + MIN_TOKEN_LIFETIME:
            return cache["token"]
    except (json.JSONDecodeError, KeyError, OSError):
        pass
    return None


def _require_token() -> str:
    token = _load_cached_token()
    if not token:
        print("TOKEN_EXPIRED", file=sys.stderr)
        sys.exit(2)
    return token


def _headers(token: str) -> dict:
    return {"Authorization": f"Bearer {token}"}


# ---------------------------------------------------------------------------
# Graph API helpers
# ---------------------------------------------------------------------------

def _search_drive_item(query: str, token: str, max_results: int = 10) -> list[dict]:
    """Sucht driveItems ueber die Graph Search API."""
    url = f"{GRAPH_BASE}/search/query"
    body = {
        "requests": [{
            "entityTypes": ["driveItem"],
            "query": {"queryString": query},
            "from": 0,
            "size": max_results,
        }]
    }
    r = requests.post(url, headers={**_headers(token), "Content-Type": "application/json"},
                      json=body, timeout=20)
    if r.status_code == 401:
        print("TOKEN_EXPIRED", file=sys.stderr)
        sys.exit(2)
    if r.status_code != 200:
        print(f"ERROR {r.status_code}: {r.text[:300]}", file=sys.stderr)
        sys.exit(1)

    data = r.json()
    hits = data["value"][0].get("hitsContainers", [{}])[0].get("hits", [])
    results = []
    for hit in hits:
        res = hit.get("resource", {})
        parent = res.get("parentReference", {})
        results.append({
            "name": res.get("name", ""),
            "webUrl": res.get("webUrl", ""),
            "itemId": hit.get("hitId", ""),
            "driveId": parent.get("driveId", ""),
            "siteId": parent.get("siteId", ""),
            "lastModified": res.get("lastModifiedDateTime", ""),
            "size": res.get("size", 0),
        })
    return results


def _resolve_url_to_drive_item(url: str, token: str) -> tuple[str, str]:
    """Loest eine SharePoint/OneDrive-URL in (driveId, itemId) auf."""
    # Methode 1: shares/encode API fuer beliebige SharePoint-URLs
    import base64
    encoded = base64.urlsafe_b64encode(url.encode()).decode().rstrip("=")
    share_url = f"{GRAPH_BASE}/shares/u!{encoded}/driveItem"
    r = requests.get(share_url, headers=_headers(token), timeout=15)
    if r.status_code == 200:
        item = r.json()
        parent = item.get("parentReference", {})
        return parent.get("driveId", ""), item.get("id", "")
    if r.status_code == 401:
        print("TOKEN_EXPIRED", file=sys.stderr)
        sys.exit(2)

    # Methode 2: Fallback — Suche nach Dateiname
    parsed = urlparse(url)
    filename = unquote(parsed.path.split("/")[-1])
    results = _search_drive_item(filename, token, max_results=5)
    # Besten Match finden
    for res in results:
        if res["name"] == filename or filename in res["webUrl"]:
            return res["driveId"], res["itemId"]

    print(f"ERROR: Konnte URL nicht aufloesen: {url}", file=sys.stderr)
    sys.exit(1)


def _download_content(drive_id: str, item_id: str, token: str) -> bytes:
    """Laedt den Datei-Inhalt als Bytes."""
    url = f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}/content"
    r = requests.get(url, headers=_headers(token), timeout=60)
    if r.status_code == 401:
        print("TOKEN_EXPIRED", file=sys.stderr)
        sys.exit(2)
    if r.status_code not in (200, 302):
        print(f"ERROR {r.status_code}: {r.text[:300]}", file=sys.stderr)
        sys.exit(1)
    return r.content


def _get_metadata(drive_id: str, item_id: str, token: str) -> dict:
    """Holt Metadaten eines driveItems."""
    url = f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}"
    r = requests.get(url, headers=_headers(token), timeout=15)
    if r.status_code == 200:
        return r.json()
    return {}


# ---------------------------------------------------------------------------
# Format-spezifische Parser
# ---------------------------------------------------------------------------

def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))
        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Notes: {notes}]")
        if texts:
            lines.append(f"\n### Folie {i}\n")
            lines.extend(texts)
    return "\n".join(lines)


def _parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"\n### Sheet: {ws.title}\n")
        lines.append("```csv")
        has_rows = False
        for row in ws.iter_rows(max_row=500, values_only=True):
            has_rows = True
            row_strs = [str(c) if c is not None else "" for c in row]
            escaped = []
            for val in row_strs:
                if ";" in val or "\n" in val or '"' in val:
                    escaped.append('"' + val.replace('"', '""') + '"')
                else:
                    escaped.append(val)
            lines.append(";".join(escaped))
        if not has_rows:
            lines.append("_(leer)_")
        lines.append("```")
    wb.close()
    return "\n".join(lines)


def _parse_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(BytesIO(data))
    lines = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style:
            lines.append(f"\n# {t}")
        elif "Heading 2" in style:
            lines.append(f"\n## {t}")
        elif "Heading 3" in style:
            lines.append(f"\n### {t}")
        else:
            lines.append(t)
    return "\n".join(lines)


def _parse_csv_content(data: bytes) -> str:
    text = data.decode("utf-8-sig", errors="replace")
    reader = csv.reader(StringIO(text), delimiter=";")
    rows = []
    for i, row in enumerate(reader):
        if i >= 200:
            rows.append(["_(... gekuerzt nach 200 Zeilen)_"])
            break
        rows.append(row)
    if not rows:
        return "_(leer)_"
    max_cols = max(len(r) for r in rows)
    for r in rows:
        r.extend([""] * (max_cols - len(r)))
    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _parse_pdf(data: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        return "_(PDF-Extraktion benoetigt 'pdfplumber': pip install pdfplumber)_"
    lines = []
    with pdfplumber.open(BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                lines.append(f"\n### Seite {i}\n")
                lines.append(text.strip())
    return "\n".join(lines) if lines else "_(Kein extrahierbarer Text im PDF)_"


def _parse_plaintext(data: bytes) -> str:
    return data.decode("utf-8-sig", errors="replace")


def _parse_image(data: bytes, name: str) -> str:
    lines = [f"_(Bilddatei: {name}, {len(data):,} bytes)_"]
    suffix = Path(name).suffix.lower()
    lines.append(f"- **Format:** {suffix}")
    try:
        from PIL import Image
        img = Image.open(BytesIO(data))
        w, h = img.size
        mode = img.mode
        lines.append(f"- **Abmessungen:** {w} x {h} px")
        lines.append(f"- **Modus:** {mode}")
    except ImportError:
        lines.append("_(Dimensionen benoetigen 'Pillow': pip install Pillow)_")
    except Exception:
        pass
    lines.append("\n_Bild wurde geladen. Nutze --download um es an einem bestimmten Ort zu speichern._")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------

def cmd_search(query: str) -> None:
    token = _require_token()
    results = _search_drive_item(query, token)
    print(f"### Graph Search: \"{query}\"\n")
    print(f"**{len(results)} Treffer**\n")
    if not results:
        print("Keine Ergebnisse.")
        return

    print("| # | Name | Groesse | Letzte Aenderung | driveId | itemId |")
    print("|---|------|---------|-------------------|---------|--------|")
    for i, r in enumerate(results, 1):
        size_kb = r["size"] // 1024 if r["size"] else "?"
        modified = r["lastModified"][:10] if r["lastModified"] else "?"
        name = r["name"]
        if len(name) > 60:
            name = name[:57] + "..."
        print(f"| {i} | [{name}]({r['webUrl']}) | {size_kb} KB | {modified} | `{r['driveId'][:20]}...` | `{r['itemId'][:20]}...` |")


def cmd_read(target: str, download_path: str | None = None) -> None:
    token = _require_token()

    # Bestimme driveId + itemId
    if target.startswith("http"):
        drive_id, item_id = _resolve_url_to_drive_item(target, token)
    elif "|" in target:
        # Format: driveId|itemId
        drive_id, item_id = target.split("|", 1)
    else:
        # Annahme: Dateiname → suchen
        results = _search_drive_item(target, token, max_results=3)
        if not results:
            print(f"ERROR: Keine Datei gefunden fuer: {target}", file=sys.stderr)
            sys.exit(1)
        drive_id = results[0]["driveId"]
        item_id = results[0]["itemId"]
        print(f"_Gefunden: {results[0]['name']}_\n", file=sys.stderr)

    # Metadaten holen
    meta = _get_metadata(drive_id, item_id, token)
    name = meta.get("name", "unbekannt")
    size = meta.get("size", 0)
    modified = meta.get("lastModifiedDateTime", "?")
    web_url = meta.get("webUrl", "")
    modified_by = meta.get("lastModifiedBy", {}).get("user", {}).get("displayName", "?")

    print(f"## {name}\n")
    print(f"- **Groesse:** {size:,} bytes ({size // 1024} KB)")
    print(f"- **Geaendert:** {modified}")
    print(f"- **Geaendert von:** {modified_by}")
    if web_url:
        print(f"- **URL:** {web_url}")
    print()

    # Datei-Inhalt laden
    data = _download_content(drive_id, item_id, token)

    # Optional lokal speichern
    if download_path:
        out = Path(download_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_bytes(data)
        print(f"_Gespeichert: {out}_\n")

    # Format-spezifisch parsen
    suffix = Path(name).suffix.lower()
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".svg", ".webp"}
    parsers = {
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".docx": _parse_docx,
        ".csv": _parse_csv_content,
        ".pdf": _parse_pdf,
        ".txt": _parse_plaintext,
        ".md": _parse_plaintext,
        ".json": _parse_plaintext,
        ".xml": _parse_plaintext,
        ".html": _parse_plaintext,
        ".htm": _parse_plaintext,
        ".log": _parse_plaintext,
    }

    parser = parsers.get(suffix)
    if parser:
        print("---\n")
        try:
            print(parser(data))
        except Exception as e:
            print(f"ERROR beim Parsen ({suffix}): {e}", file=sys.stderr)
            print(f"_(Datei wurde geladen ({len(data)} bytes), aber das Parsen schlug fehl.)_")
    elif suffix in image_exts:
        if not download_path:
            auto_path = Path(__file__).resolve().parent.parent / "userdata" / "tmp" / name
            auto_path.parent.mkdir(parents=True, exist_ok=True)
            auto_path.write_bytes(data)
            print(f"_Bild gespeichert: {auto_path}_\n")
        print("---\n")
        print(_parse_image(data, name))
    else:
        print(f"_(Format '{suffix}' wird nicht als Text extrahiert. Nutze --download um die Datei lokal zu speichern.)_")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1]

    if cmd == "search":
        if len(sys.argv) < 3:
            print("Usage: python m365_file_reader.py search \"Dateiname\"", file=sys.stderr)
            sys.exit(1)
        cmd_search(sys.argv[2])

    elif cmd == "read":
        if len(sys.argv) < 3:
            print("Usage: python m365_file_reader.py read URL_ODER_PFAD [--download ZIEL]", file=sys.stderr)
            sys.exit(1)
        target = sys.argv[2]
        download_path = None
        if "--download" in sys.argv:
            idx = sys.argv.index("--download")
            if idx + 1 < len(sys.argv):
                download_path = sys.argv[idx + 1]
        cmd_read(target, download_path)

    else:
        print(f"Unbekannter Befehl: {cmd}", file=sys.stderr)
        print(__doc__)
        sys.exit(1)


if __name__ == "__main__":
    main()
