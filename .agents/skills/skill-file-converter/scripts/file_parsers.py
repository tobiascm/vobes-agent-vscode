"""Gemeinsame Format-Parser fuer M365-Dateien und Mail-Anhaenge.

Konvertiert Bytes in lesbaren Text. Wird von m365_file_reader.py
und m365_mail_search.py verwendet.

Unterstuetzte Formate:
    .pptx  — Text pro Folie (inkl. Tabellen, Notes)
    .xlsx  — Alle Sheets als CSV (Semikolon-getrennt)
    .docx  — Volltext mit Ueberschriften
    .pdf   — Text-Extraktion (pdfplumber)
    .csv   — Erste 200 Zeilen als Markdown-Tabelle
    .txt/.md/.json/.xml — Plaintext
    .png/.jpg/... — Bild-Metadaten
"""

from __future__ import annotations

import csv
from io import BytesIO, StringIO
from pathlib import Path


def parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Notes: {notes}]")
        if texts:
            lines.append(f"\n### Folie {i}\n")
            lines.extend(texts)
    return "\n".join(lines)


def parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"\n### Sheet: {ws.title}\n")
        lines.append("```csv")
        has_rows = False
        for row in ws.iter_rows(max_row=500, values_only=True):
            has_rows = True
            row_strs = [str(c) if c is not None else "" for c in row]
            escaped = []
            for val in row_strs:
                if ";" in val or "\n" in val or '"' in val:
                    escaped.append('"' + val.replace('"', '""') + '"')
                else:
                    escaped.append(val)
            lines.append(";".join(escaped))
        if not has_rows:
            lines.append("_(leer)_")
        lines.append("```")
    wb.close()
    return "\n".join(lines)


def parse_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(BytesIO(data))
    lines = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style:
            lines.append(f"\n# {t}")
        elif "Heading 2" in style:
            lines.append(f"\n## {t}")
        elif "Heading 3" in style:
            lines.append(f"\n### {t}")
        else:
            lines.append(t)
    return "\n".join(lines)


def parse_csv_content(data: bytes) -> str:
    text = data.decode("utf-8-sig", errors="replace")
    reader = csv.reader(StringIO(text), delimiter=";")
    rows = []
    for i, row in enumerate(reader):
        if i >= 200:
            rows.append(["_(... gekuerzt nach 200 Zeilen)_"])
            break
        rows.append(row)
    if not rows:
        return "_(leer)_"
    max_cols = max(len(r) for r in rows)
    for r in rows:
        r.extend([""] * (max_cols - len(r)))
    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def parse_pdf(data: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        return "_(PDF-Extraktion benoetigt 'pdfplumber': pip install pdfplumber)_"
    lines = []
    with pdfplumber.open(BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                lines.append(f"\n### Seite {i}\n")
                lines.append(text.strip())
    return "\n".join(lines) if lines else "_(Kein extrahierbarer Text im PDF)_"


def parse_plaintext(data: bytes) -> str:
    return data.decode("utf-8-sig", errors="replace")


def parse_image(data: bytes, name: str) -> str:
    lines = [f"_(Bilddatei: {name}, {len(data):,} bytes)_"]
    suffix = Path(name).suffix.lower()
    lines.append(f"- **Format:** {suffix}")
    try:
        from PIL import Image
        img = Image.open(BytesIO(data))
        w, h = img.size
        mode = img.mode
        lines.append(f"- **Abmessungen:** {w} x {h} px")
        lines.append(f"- **Modus:** {mode}")
    except ImportError:
        lines.append("_(Dimensionen benoetigen 'Pillow': pip install Pillow)_")
    except Exception:
        pass
    return "\n".join(lines)


# Zuordnung Dateiendung → Parser
PARSERS: dict[str, callable] = {
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
    ".docx": parse_docx,
    ".csv": parse_csv_content,
    ".pdf": parse_pdf,
    ".txt": parse_plaintext,
    ".md": parse_plaintext,
    ".json": parse_plaintext,
    ".xml": parse_plaintext,
    ".html": parse_plaintext,
    ".htm": parse_plaintext,
    ".log": parse_plaintext,
}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".svg", ".webp"}


def convert_bytes(data: bytes, filename: str) -> str | None:
    """Konvertiert Datei-Bytes in Text. Gibt None zurueck wenn Format nicht unterstuetzt."""
    suffix = Path(filename).suffix.lower()
    parser = PARSERS.get(suffix)
    if parser:
        return parser(data)
    if suffix in IMAGE_EXTS:
        return parse_image(data, filename)
    return None
